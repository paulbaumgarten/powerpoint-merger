#!/usr/bin/env python3.7

# Powerpoint building/editing class
# Created by: Paul Baumgarten 2019
# Last updated: 08.06.2020 by Paul Baumgarten

# https://python-pptx.readthedocs.io

import pptx
from string import Template
import copy
import json
import os
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor

class PowerPointer():
    def __init__(self, ppt_filename, media_folder="."):
        self.ppt_filename = ppt_filename
        self.ppt = pptx.Presentation(ppt_filename)
        self.media_folder = media_folder
    
    def save(self, save_ppt_filename):
        """ Saves PPT to the local file system """
        self.ppt.save(save_ppt_filename)

    def get_all_slide_ids(self):
        """ Returns list of all slide_id integers, in order of appearance in presentation """
        coll = []
        for i in range(len(self.ppt.slides)):
            coll.append( self.ppt.slides[i].slide_id )
        return coll

    def get_slide_id_from_slide_number(self, slide_number):
        """ Returns slide_id integer """
        return self.ppt.slides[slide_number].slide_id

    def get_slide_by_id(self, slide_id):
        """ Returns Slide object """
        return self.ppt.slides.get(slide_id)

    def get_slide_by_number(self, slide_number):
        """ Returns Slide object """
        if len(self.ppt.slides) > slide_number:
            return self.ppt.slides[slide_number]
        else:
            return None
    
    def get_slide_elements(self, slide_id):
        shapes = self.get_slide_by_id(slide_id).shapes
        return shapes
    
    def copy_slide(self, from_slide_id):
        """
        Duplicate the slide with the given from_slide_id in pres.
        Appends the new slide to the end of the presentation
        Returns Slide object
        """
        # from https://github.com/scanny/python-pptx/issues/132#issuecomment-346699019
        def _get_blank_slide_layout():
            layout_items_count = [len(layout.placeholders) for layout in self.ppt.slide_layouts]
            min_items = min(layout_items_count)
            blank_layout_id = layout_items_count.index(min_items)
            return self.ppt.slide_layouts[blank_layout_id]

        source = self.ppt.slides.get(from_slide_id)
        blank_slide_layout = _get_blank_slide_layout()
        dest = self.ppt.slides.add_slide(blank_slide_layout)

        for shape in source.shapes:
            newel = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for key, value in source.part.rels.items():
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
        return dest

    def new_slide(self, layout_name="Blank"):
        """
        Creates new slide from the given layout slide
        Appends the new slide to the end of the presentation
        Returns Slide object

        To find a layout name, in PPT:
         -> edit master slides
         -> right click on the thumbnail of the layout you want
         -> select `rename layout`
        """
        slide_layout = self.ppt.slide_layouts.get_by_name(layout_name)
        slide = self.ppt.slides.add_slide(slide_layout)
        return slide

    def parse_placeholders(self, slide_id, recordset):
        slide = self.ppt.slides.get(slide_id)
         
        # Get info from layout about the placeholders in this slide
        layout_placeholders = {}
        for lph in slide.slide_layout.placeholders:
            # --> .placeholder_format.idx is the lookup index to use for referencing the layout placeholder later
            # --> refer https://github.com/scanny/python-pptx/issues/475
            rec = {"shape_id": lph.shape_id, "shape_idx": lph.placeholder_format.idx, "shape_type": lph.shape_type, "text": lph.text}
            layout_placeholders[ lph.placeholder_format.idx ] = rec
        # print(json.dumps(layout_placeholders, indent=3)) # for debugging help

        for shape in slide.placeholders:
            layout_placeholder_index = shape.placeholder_format.idx
            placeholder_text = layout_placeholders[layout_placeholder_index]["text"]
            
            # Now we have our placeholder text, perform the subsitution lookup with the provided data recordset
            try:
                rendered = Template(placeholder_text).substitute(recordset)
            except:
                rendered = "MISSING FIELD ("+placeholder_text+")"
                print(f"[parse_placeholders] Slide_ID {slide_id} missing field `{placeholder_text}`")

            # Substitute our rendered text into the shape
            if shape.placeholder_format.type == 18: # Picture
                # the text should turn into a filename
                filename = os.path.join(self.media_folder, rendered)
                try:
                    shape.insert_picture(filename)
                except FileNotFoundError:
                    print(f"[parse_placeholders] Slide_ID {slide_id} image file `{rendered}` file not found")
                except ValueError:
                    print(f"[parse_placeholders] Slide_ID {slide_id} image file `{rendered}` unsupported image format")
                except:
                    print(f"[parse_placeholders] Slide_ID {slide_id} image file `{rendered}` unknown error with image")
            else: # Assume it's text for now
                if shape.has_text_frame:
                    shape.text = rendered

    def parse_slide_content(self, slide_id, recordset):
        shapes = self.ppt.slides.get(slide_id).shapes
        for shape in shapes:
            if shape.has_text_frame:
                try:
                    rendered = Template(shape.text).substitute(recordset)
                except:
                    rendered = "MISSING FIELD ("+shape.text+")"
                shape.text = rendered

    def add_image(self, slide_id, file, x, y, h, w, **kwargs):
        slide = self.ppt.slides.get(slide_id)
        left = Mm(x)
        top = Mm(y)
        height = Mm(h)
        width = Mm(w)
        try:
            slide.shapes.add_picture(file, left, top, height=height, width=width)
        except FileNotFoundError:
            print(f"[Filenotfound] in Powerpointer.add_image() for file `{file}`")

    def add_text(self, slide_id, text, x, y, w, h=1, font_name=None, color=None, bold=False, size=18, **kwargs):
        def hex_to_rgb(value):
            """ Converts `#ffffff` to (255,255,255) """
            value = value.lstrip('#')
            r = int(value[:2], 16)
            g = int(value[2:4], 16)
            b = int(value[4:], 16)
            return (r,g,b)

        slide = self.ppt.slides.get(slide_id)
        left = Mm(x)
        top = Mm(y)
        width = Mm(w)
        height = Mm(h)
        tx_box = slide.shapes.add_textbox(left, top, width, height)     # Add a text box to the slide
        tx_frame = tx_box.text_frame                                    # Add a text frame to the text box
        p = tx_frame.add_paragraph()                                    # Add a paragraph to the text frame
        p.text = text
        # Apply formatting options
        p.font.bold = bold
        p.font.size = Pt(size)
        if color is not None:
            p.font.color.rgb = RGBColor(*hex_to_rgb(color))
        if font_name is not None:
            p.font.name = font_name


"""
    Works as of 25.01.2019 but was based off a workaround that I discovered through experimentation. The new version is based off a github support comment from the developer.

    def parse_placeholders_(self, slide_id, recordset):
        slide = self.ppt.slides.get(slide_id)
        #print("")
        #print(f"[parse_placeholders] slide_id {slide_id}")
         
        # Find info from layout about placeholders in this slide
        layout_placeholders = []
        for lph in slide.slide_layout.placeholders:
            rec = {"shape_id": lph.shape_id, "shape_type": lph.shape_type, "text": lph.text}
            layout_placeholders.append(rec)

        for shape in slide.placeholders:
            #print("[parse_placeholders] shape_id: ",shape.shape_id)
            #print("[parse_placeholders] shape_type: ",shape.shape_type)
            #print("[parse_placeholders] text: ",shape.text)
            #print("[parse_placeholders] slide layout name: ",shape.part.slide_layout.name)
            #print("[parse_placeholders] name: ",shape.name)

            layout_placeholder_index = int(shape.name[shape.name.rfind(" ")+1:])-1 # don't ask...
            # ok, since you asked for it... it was a nightmare attempting to figure out the
            # layout shape that corresponded to the shapes created on a slide that was made based on a layout.
            # in the end, it *seems* as though, the shape.name ends with an integer that is increments through the slide.
            # that integer also seems to correlate to the order in which the slide_layout.placeholders appear.
            # so, if it works, i'm holding to that assumption until it breaks.... i'm not hopeful.

            #print("[parse_placeholders] layout_placeholder_index: ",layout_placeholder_index)
            placeholder_text = layout_placeholders[layout_placeholder_index]["text"]
            #print("[parse_placeholders] placeholder_text: ",placeholder_text)
            
            # Now we have our placeholder text, perform the subsitution lookup with the provided data recordset
            try:
                rendered = Template(placeholder_text).substitute(recordset)
            except:
                rendered = "MISSING FIELD ("+placeholder_text+")"

            # Substitute our rendered text into the shape
            if shape.placeholder_format.type == 18: # Picture
                # the text should turn into a filename
                try:
                    shape.insert_picture(rendered)
                except FileNotFoundError:
                    print("[parse_placeholders] Image file not found: "+rendered)
            else: # Assume it's text for now
                if shape.has_text_frame:
                    shape.text = rendered
"""